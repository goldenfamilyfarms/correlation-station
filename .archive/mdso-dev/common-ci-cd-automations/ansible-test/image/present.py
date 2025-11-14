import os
from pptx import Presentation

prs = Presentation()
title_slide_layout = prs.slide_layouts[0]
prs.slide_width = 10076688
prs.slide_height = 5669280

entries = os.listdir('slides/')
for entry in sorted(entries):
    slide = prs.slides.add_slide(title_slide_layout)
    shapes = slide.shapes
    print(entry)
    picture = shapes.add_picture('slides/' + entry, 0, 0, height=5669280, width=10076688)

prs.save('presentation.pptx')
